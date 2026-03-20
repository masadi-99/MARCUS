"""
Create a professional animated PPTX presentation showcasing the MARCUS pipeline.

Designed for short-form social media video (LinkedIn, Twitter/X).
Style: minimalistic, elegant, academic, dark navy + white + blue accent.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Colors ──────────────────────────────────────────────────────────────
NAVY = RGBColor(0x0F, 0x17, 0x2A)         # Dark navy background
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)    # Primary accent
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)  # Lighter blue accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)   # Subtle text
WARM_WHITE = RGBColor(0xF8, 0xFA, 0xFC)   # Off-white
GREEN = RGBColor(0x22, 0xC5, 0x5E)        # Success/confidence
AMBER = RGBColor(0xF5, 0x9E, 0x0B)        # Warning

# ── Slide dimensions (16:9) ─────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Inter", italic=False, spacing=None):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    if spacing is not None:
        p.space_after = Pt(spacing)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     border_color=None, text="", font_size=14,
                     text_color=WHITE, bold=True):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    # Adjust corner radius
    shape.adjustments[0] = 0.15

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = bold
        p.font.name = "Inter"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top, color=ACCENT_BLUE, width=Pt(2)):
    """Add a connector arrow."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    connector.line._element.attrib[qn('a:tailEnd')] = ''
    line_elem = connector.line._element
    tail = line_elem.makeelement(qn('a:tailEnd'), {})
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    line_elem.append(tail)
    return connector


def add_entrance_animation(slide, shape, delay_ms=0, duration_ms=500, anim_type="appear"):
    """Add entrance animation to a shape.

    Note: python-pptx doesn't have full animation support, so we use
    raw XML manipulation to add PowerPoint-compatible animations.
    """
    # Get the shape's spTree (parent) element
    sp_tree = slide.shapes._spTree

    # Create timing node if it doesn't exist
    sld = slide._element
    timing = sld.find(qn('p:timing'))
    if timing is None:
        timing = sld.makeelement(qn('p:timing'), {})
        sld.append(timing)

    tn_lst = timing.find(qn('p:tnLst'))
    if tn_lst is None:
        tn_lst = timing.makeelement(qn('p:tnLst'), {})
        timing.append(tn_lst)

    par_seq = tn_lst.find(qn('p:par'))
    if par_seq is None:
        par_seq = tn_lst.makeelement(qn('p:par'), {})
        tn_lst.append(par_seq)

        c_tn = par_seq.makeelement(qn('p:cTn'), {
            'id': '1', 'dur': 'indefinite', 'restart': 'never', 'nodeType': 'tmRoot'
        })
        par_seq.append(c_tn)

        child_tn_lst = c_tn.makeelement(qn('p:childTnLst'), {})
        c_tn.append(child_tn_lst)

        seq = child_tn_lst.makeelement(qn('p:seq'), {
            'concurrent': '1', 'nextAc': 'seek'
        })
        child_tn_lst.append(seq)

        seq_ctn = seq.makeelement(qn('p:cTn'), {
            'id': '2', 'dur': 'indefinite', 'nodeType': 'mainSeq'
        })
        seq.append(seq_ctn)

        seq_child_lst = seq_ctn.makeelement(qn('p:childTnLst'), {})
        seq_ctn.append(seq_child_lst)

        # prev condition
        prev_cond = seq.makeelement(qn('p:prevCondLst'), {})
        seq.append(prev_cond)
        cond = prev_cond.makeelement(qn('p:cond'), {'evt': 'onPrev', 'delay': '0'})
        prev_cond.append(cond)
        tgt_el = cond.makeelement(qn('p:tgtEl'), {})
        cond.append(tgt_el)
        sld_tgt = tgt_el.makeelement(qn('p:sldTgt'), {})
        tgt_el.append(sld_tgt)

        # next condition
        next_cond = seq.makeelement(qn('p:nextCondLst'), {})
        seq.append(next_cond)
        cond2 = next_cond.makeelement(qn('p:cond'), {'evt': 'onNext', 'delay': '0'})
        next_cond.append(cond2)
        tgt_el2 = cond2.makeelement(qn('p:tgtEl'), {})
        cond2.append(tgt_el2)
        sld_tgt2 = tgt_el2.makeelement(qn('p:sldTgt'), {})
        tgt_el2.append(sld_tgt2)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY)

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
                 "MARCUS", font_size=72, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Inter")

    # Subtitle
    add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(0.8),
                 "Multimodal Autonomous Reasoning and Chat\nfor Ultrasound and Signals",
                 font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
                 font_name="Inter", italic=True)

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.67), Inches(3.3), Inches(2), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Authors
    add_text_box(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.6),
                 "O'Sullivan JW*, Asadi M*, Elbe L, Chaudhari A, Nedaee T,\n"
                 "Haddad F, Salerno M, Fei-Fei L, Adeli E, Arnaout R, Ashley EA",
                 font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Affiliation
    add_text_box(slide, Inches(2), Inches(5.6), Inches(9), Inches(0.4),
                 "STANFORD UNIVERSITY  |  UCSF",
                 font_size=11, color=RGBColor(0x64, 0x74, 0x8B),
                 alignment=PP_ALIGN.CENTER, font_name="Inter")

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 2: The Problem — three cardiac modalities
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
                 "THREE MODALITIES, ONE PATIENT",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT,
                 font_name="Inter")

    # Three modality cards with images
    card_w = Inches(3.4)
    card_h = Inches(4.5)
    card_y = Inches(1.5)
    gap = Inches(0.6)
    start_x = Inches(1.3)

    modalities = [
        ("ECG", "12-lead Electrocardiogram", "/home/masadi/temp_ecg.png"),
        ("Echo", "Echocardiogram", "/tmp/echo_frame.png"),
        ("CMR", "Cardiac MRI", "/tmp/cmr_frame.png"),
    ]

    for i, (name, desc, img_path) in enumerate(modalities):
        x = start_x + i * (card_w + gap)

        # Card background
        card = add_rounded_rect(slide, x, card_y, card_w, card_h,
                                RGBColor(0x1A, 0x24, 0x3B),
                                border_color=RGBColor(0x2D, 0x3A, 0x55))

        # Modality label
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.25),
                     card_w - Inches(0.6), Inches(0.4),
                     name, font_size=24, color=WHITE, bold=True,
                     alignment=PP_ALIGN.LEFT)

        # Description
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.7),
                     card_w - Inches(0.6), Inches(0.3),
                     desc, font_size=11, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.LEFT)

        # Image
        if os.path.isfile(img_path):
            try:
                slide.shapes.add_picture(
                    img_path,
                    x + Inches(0.3), card_y + Inches(1.2),
                    card_w - Inches(0.6), Inches(2.8)
                )
            except Exception:
                pass

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 3: Training Scale
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
                 "TRAINING AT SCALE",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # Big numbers
    stats = [
        ("13.5M", "Images"),
        ("270K", "Clinical Studies"),
        ("3", "Expert Models"),
    ]

    stat_w = Inches(3.2)
    stat_start = Inches(1.3)
    stat_y = Inches(2.3)

    for i, (val, label) in enumerate(stats):
        x = stat_start + i * (stat_w + Inches(0.4))

        add_text_box(slide, x, stat_y, stat_w, Inches(1.2),
                     val, font_size=64, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="JetBrains Mono")

        add_text_box(slide, x, stat_y + Inches(1.3), stat_w, Inches(0.5),
                     label, font_size=16, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(1),
                 "Physician-verified ground truth across ECG, echocardiography, and cardiac MRI.\n"
                 "Group-relative policy optimization (GRPO) and supervised fine-tuning (SFT).",
                 font_size=14, color=RGBColor(0x64, 0x74, 0x8B),
                 alignment=PP_ALIGN.CENTER, italic=True)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 4: Architecture — Expert Routing
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
                 "AGENTIC ARCHITECTURE",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # Input box (left)
    add_rounded_rect(slide, Inches(0.6), Inches(2.5), Inches(2.2), Inches(2.5),
                     RGBColor(0x1A, 0x24, 0x3B),
                     border_color=RGBColor(0x3B, 0x82, 0xF6),
                     text="Patient\nData", font_size=18, text_color=WHITE)

    # Arrow from input to experts
    arrow_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2.9), Inches(3.7), Inches(0.8), Pt(2)
    )
    arrow_line.fill.solid()
    arrow_line.fill.fore_color.rgb = ACCENT_BLUE
    arrow_line.line.fill.background()

    # Expert boxes (middle column)
    expert_x = Inches(3.9)
    expert_w = Inches(2.2)
    expert_h = Inches(0.7)

    experts = [
        ("ECG Expert", RGBColor(0x16, 0x65, 0x34)),
        ("Echo Expert", RGBColor(0x1E, 0x3A, 0x5F)),
        ("CMR Expert", RGBColor(0x6D, 0x28, 0xD9)),
    ]

    for i, (name, color) in enumerate(experts):
        y = Inches(1.8) + i * Inches(1.1)
        add_rounded_rect(slide, expert_x, y, expert_w, expert_h,
                         color, text=name, font_size=14)

    # Arrow from experts to mirage probe
    arrow_line2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.3), Inches(3.7), Inches(0.8), Pt(2)
    )
    arrow_line2.fill.solid()
    arrow_line2.fill.fore_color.rgb = ACCENT_BLUE
    arrow_line2.line.fill.background()

    # Mirage Probe box
    add_rounded_rect(slide, Inches(7.3), Inches(2.5), Inches(2.2), Inches(2.5),
                     RGBColor(0x78, 0x35, 0x0F),
                     border_color=AMBER,
                     text="Mirage\nProbe", font_size=18, text_color=WHITE)

    # Arrow from mirage to orchestrator
    arrow_line3 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.7), Inches(3.7), Inches(0.8), Pt(2)
    )
    arrow_line3.fill.solid()
    arrow_line3.fill.fore_color.rgb = ACCENT_BLUE
    arrow_line3.line.fill.background()

    # Orchestrator + Report box
    add_rounded_rect(slide, Inches(10.7), Inches(2.5), Inches(2.2), Inches(2.5),
                     DARK_BLUE,
                     border_color=ACCENT_BLUE,
                     text="Orchestrator\n& Synthesis", font_size=18, text_color=WHITE)

    # Labels below boxes
    labels = [
        (Inches(0.6), "Input"),
        (Inches(3.9), "Specialized Experts"),
        (Inches(7.3), "Hallucination Guard"),
        (Inches(10.7), "Unified Report"),
    ]
    for x, txt in labels:
        add_text_box(slide, x, Inches(5.3), Inches(2.2), Inches(0.4),
                     txt, font_size=10, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 5: Expert Analysis — showing attention
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
                 "EXPERT ANALYSIS WITH ATTENTION",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # Show ECG with attention heatmap
    ecg_path = "/home/masadi/temp_ecg.png"
    heatmap_path = "/tmp/test_heatmap_real.png"

    if os.path.isfile(ecg_path):
        # Original
        add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
                     "Input: 12-lead ECG", font_size=12, color=LIGHT_GRAY)
        slide.shapes.add_picture(ecg_path, Inches(0.8), Inches(1.8),
                                 Inches(5.5), Inches(4.2))

    if os.path.isfile(heatmap_path):
        # With attention overlay
        add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.4),
                     "Attention Map: Where the model looks", font_size=12, color=LIGHT_GRAY)
        slide.shapes.add_picture(heatmap_path, Inches(7), Inches(1.8),
                                 Inches(5.5), Inches(4.2))

    # Arrow between
    arrow_mid = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(6.2), Inches(3.5), Inches(0.6), Inches(0.5)
    )
    arrow_mid.fill.solid()
    arrow_mid.fill.fore_color.rgb = ACCENT_BLUE
    arrow_mid.line.fill.background()

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 6: Mirage Detection
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
                 "MIRAGE DETECTION",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
                 "Do vision-language models actually see, or just fabricate?",
                 font_size=28, color=WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # Left: with image column
    col_w = Inches(5)
    add_text_box(slide, Inches(0.8), Inches(2.4), col_w, Inches(0.4),
                 "WITH IMAGE", font_size=11, color=GREEN,
                 bold=True)

    queries = ["Original question", "Rephrase 1", "Rephrase 2"]
    for i, q in enumerate(queries):
        y = Inches(3.0) + i * Inches(0.8)
        add_rounded_rect(slide, Inches(0.8), y, Inches(4.8), Inches(0.6),
                         RGBColor(0x1A, 0x24, 0x3B),
                         border_color=RGBColor(0x22, 0xC5, 0x5E),
                         text=f"  {q}  →  Consistent, specific response",
                         font_size=11, text_color=WARM_WHITE, bold=False)

    # Right: without image
    add_text_box(slide, Inches(7), Inches(2.4), col_w, Inches(0.4),
                 "WITHOUT IMAGE (COUNTERFACTUAL)", font_size=11,
                 color=AMBER, bold=True)

    add_rounded_rect(slide, Inches(7), Inches(3.0), Inches(4.8), Inches(0.6),
                     RGBColor(0x1A, 0x24, 0x3B),
                     border_color=RGBColor(0xF5, 0x9E, 0x0B),
                     text="  Same question  →  Generic, vague response",
                     font_size=11, text_color=WARM_WHITE, bold=False)

    # Comparison arrow
    add_text_box(slide, Inches(5), Inches(4.3), Inches(3), Inches(0.5),
                 "Compare responses", font_size=13, color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER, bold=True)

    # Result
    add_rounded_rect(slide, Inches(3.5), Inches(5.2), Inches(6), Inches(1.2),
                     RGBColor(0x0A, 0x1E, 0x0A),
                     border_color=GREEN,
                     text="High divergence = Grounded analysis (no mirage)\n"
                          "MARCUS: 0% mirage rate vs GPT-4o: 38% vs Gemini: 35%",
                     font_size=13, text_color=WHITE, bold=False)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 7: Multimodal Orchestration
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
                 "MULTIMODAL ORCHESTRATION",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # Three expert responses flowing into orchestrator
    expert_data = [
        ("ECG", "Atrial fibrillation\nwith rapid ventricular\nresponse", RGBColor(0x16, 0x65, 0x34)),
        ("Echo", "Reduced LV function\nEF 35%, moderate\nmitral regurgitation", RGBColor(0x1E, 0x3A, 0x5F)),
        ("CMR", "Late gadolinium\nenhancement in\ninferolateral wall", RGBColor(0x6D, 0x28, 0xD9)),
    ]

    for i, (name, finding, color) in enumerate(expert_data):
        x = Inches(0.5) + i * Inches(3.2)
        y = Inches(1.5)

        # Expert label
        add_text_box(slide, x, y, Inches(2.8), Inches(0.4),
                     name, font_size=16, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Finding box
        add_rounded_rect(slide, x, y + Inches(0.5), Inches(2.8), Inches(1.8),
                         RGBColor(0x1A, 0x24, 0x3B),
                         border_color=color,
                         text=finding, font_size=12, text_color=WARM_WHITE, bold=False)

    # Convergence arrows
    for i in range(3):
        x = Inches(1.9) + i * Inches(3.2)
        arrow_down = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, x, Inches(4.0), Inches(0.3), Inches(0.5)
        )
        arrow_down.fill.solid()
        arrow_down.fill.fore_color.rgb = ACCENT_BLUE
        arrow_down.line.fill.background()

    # Orchestrator box
    add_rounded_rect(slide, Inches(2.5), Inches(4.7), Inches(8), Inches(2),
                     DARK_BLUE,
                     border_color=ACCENT_BLUE,
                     text="Synthesized Report\n\n"
                          "\"The patient presents with atrial fibrillation, reduced left ventricular\n"
                          "function (EF 35%), and evidence of prior myocardial infarction on CMR.\n"
                          "These findings together suggest ischemic cardiomyopathy.\"",
                     font_size=13, text_color=WHITE, bold=False)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 8: Architecture figure from the paper
    # ════════════════════════════════════════════════════════════════════
    fig1_path = "/home/masadi/MARCUS/docs/figures/fig1_architecture.png"
    if os.path.isfile(fig1_path):
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)

        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                     "ARCHITECTURE AND TRAINING PIPELINE",
                     font_size=14, color=NAVY, alignment=PP_ALIGN.LEFT)

        slide.shapes.add_picture(fig1_path, Inches(0.5), Inches(1.2),
                                 Inches(12.3), Inches(5.8))

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 9: Performance figure
    # ════════════════════════════════════════════════════════════════════
    fig2_path = "/home/masadi/MARCUS/docs/figures/fig2_performance.png"
    if os.path.isfile(fig2_path):
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)

        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                     "PERFORMANCE",
                     font_size=14, color=NAVY, alignment=PP_ALIGN.LEFT)

        slide.shapes.add_picture(fig2_path, Inches(0.5), Inches(1.2),
                                 Inches(12.3), Inches(5.8))

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 10: Closing
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
                 "MARCUS", font_size=72, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Inter")

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.67), Inches(3.2), Inches(2), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_text_box(slide, Inches(2), Inches(3.6), Inches(9), Inches(0.8),
                 "An agentic, multimodal vision-language model\n"
                 "for cardiac diagnosis and management",
                 font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
                 italic=True)

    add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.6),
                 "github.com/masadi-99/MARCUS",
                 font_size=14, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(5.6), Inches(9), Inches(0.4),
                 "Stanford University  |  UCSF",
                 font_size=12, color=RGBColor(0x64, 0x74, 0x8B),
                 alignment=PP_ALIGN.CENTER)

    # ── Save ────────────────────────────────────────────────────────────
    out_path = "/home/masadi/MARCUS/MARCUS_Pipeline_Demo.pptx"
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    return out_path


if __name__ == "__main__":
    create_presentation()
